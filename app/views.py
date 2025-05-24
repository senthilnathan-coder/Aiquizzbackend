from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from django.shortcuts import render
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import google.generativeai as genai
import base64
import tempfile
import os
import json
import razorpay
from django.conf import settings
from app.models import *
from mongoengine.errors import DoesNotExist, ValidationError
from datetime import datetime, timedelta
import cv2
import time
import hmac
import hashlib
import requests
from bs4 import BeautifulSoup
from faster_whisper import WhisperModel
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from app.serializers import *
# from rest_framework.permissions import IsAuthenticated

genai.configure(api_key=settings.GEMINI_API_KEY)
whisper_model = WhisperModel("tiny", compute_type="float32",device="cpu")
# whisper_model = whisper.load_model("base")

def extract_url_text(url):
    try:
        # Send a GET request to the URL
        response = requests.get(url)
        response.raise_for_status()  # Raise an exception for bad status codes
        
        # Parse the HTML content
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text content
        text = soup.get_text(separator='\n')
        
        # Clean up the text
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        text = '\n'.join(lines)
        
        if not text.strip():
            raise ValueError("No text content found in URL")
        return text
    except Exception as e:
        raise Exception(f"Error extracting URL text: {str(e)}")

def extract_frame(video_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as tmp:
        tmp.write(video_file.read())
        tmp_path = tmp.name

    cap = cv2.VideoCapture(tmp_path)
    frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    middle_frame = frame_count // 2
    cap.set(cv2.CAP_PROP_POS_FRAMES, middle_frame)

    success, frame = cap.read()
    frame_path = tmp_path + "_frame.jpg"

    if success:
        cv2.imwrite(frame_path, frame)

    cap.release()
    os.remove(tmp_path)

    with open(frame_path, "rb") as f:
        image_data = f.read()

    os.remove(frame_path)
    return image_data, "image/jpeg"
def transcribe_audio(audio_file):
    try:
        # Get the file extension from the original file
        file_extension = os.path.splitext(audio_file.name)[1].lower()
        if not file_extension:
            file_extension = '.mp3'  # Default to mp3 if no extension

        # Create temp file with correct extension
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp:
            # Read in chunks to handle large files
            for chunk in audio_file.chunks():
                tmp.write(chunk)
            tmp_path = tmp.name

        try:
            # Transcribe with error handling
            segments, info = whisper_model.transcribe(
                tmp_path,
                beam_size=5,
                language=None,  # Auto-detect language
                initial_prompt="This is a transcription of an audio file."
            )
            
            if not segments:
                raise ValueError("No speech detected in audio file")
            
            # Combine all segments with proper spacing
            text = ""
            for segment in segments:
                text += segment.text + " "
            
            text = text.strip()
            if not text:
                raise ValueError("Transcription produced empty result")
            
            return text

        except Exception as e:
            raise Exception(f"Transcription failed: {str(e)}")

        finally:
            # Clean up temp file
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    except Exception as e:
        raise Exception(f"Audio processing failed: {str(e)}")

def extract_word_text(word_file):
    try:
        doc = Document(word_file)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        if not text.strip():
            raise ValueError("No text content found in Word document")
        return text
    except Exception as e:
        raise Exception(f"Error extracting Word document text: {str(e)}")

def extract_pdf_text(pdf_file):
    try:
        reader = PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        if not text.strip():
            raise ValueError("No text content found in PDF document")
        return text
    except Exception as e:
        raise Exception(f"Error extracting PDF text: {str(e)}")

def extract_ppt_text(ppt_file):
    try:
        prs = Presentation(ppt_file)
        text = ""
        
        # Extract text from all slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        if not text.strip():
            raise ValueError("No text content found in PowerPoint document")
        return text
    except Exception as e:
        raise Exception(f"Error extracting PowerPoint text: {str(e)}")

def extract_excel_text(excel_file):
    try:
        workbook = load_workbook(excel_file)
        text = ""
        
        # Extract text from all sheets
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        text += str(cell.value) + "\n"
        
        if not text.strip():
            raise ValueError("No text content found in Excel document")
        return text
    except Exception as e:
        raise Exception(f"Error extracting Excel document text: {str(e)}")

def parse_questions(response_text, question_type='mcq'):
    questions = []
    blocks = response_text.strip().split("Q")[1:]
    
    try:
        for block in blocks:
            try:
                lines = [line.strip() for line in block.strip().splitlines() if line.strip()]
                
                if not lines:
                    continue
                
                # Extract question text more flexibly
                question_text = lines[0]
                if ':' in question_text:
                    question_text = question_text.split(':', 1)[1].strip()
                
                # Handle MCQ format
                if question_type.lower() == 'mcq':
                    options = []
                    answer_line = None
                    
                    # Extract options and answer more flexibly
                    for line in lines[1:]:
                        line = line.strip()
                        if line.lower().startswith(('a.', 'b.', 'c.', 'd.')):
                            options.append(line[2:].strip())
                        elif any(line.lower().startswith(prefix) for prefix in ['answer:', 'ans:', 'a:']):
                            answer_line = line
                    
                    if len(options) == 4 and answer_line:
                        correct_letter = answer_line.split(':')[1].strip().upper()[0]
                        correct_index = ord(correct_letter) - ord('A')
                        
                        if 0 <= correct_index < len(options):
                            questions.append({
                                'question': question_text,
                                'options': options,
                                'answer': options[correct_index]
                            })
                
                # Handle True/False format
                elif question_type.lower() == 'true_false':
                    options = ['True', 'False']
                    answer_line = None
                    
                    # Find answer line with flexible matching
                    for line in lines:
                        if any(line.lower().startswith(prefix) for prefix in ['answer:', 'ans:', 'a:']):
                            answer_line = line
                            break
                    
                    if answer_line:
                        answer_text = answer_line.split(':')[1].strip().lower()
                        
                        # Simplified answer matching for true/false
                        if answer_text in ['true', 't', 'a']:
                            correct_index = 0
                        elif answer_text in ['false', 'f', 'b']:
                            correct_index = 1
                        else:
                            continue
                        
                        questions.append({
                            'question': question_text,
                            'options': options,
                            'answer': options[correct_index]
                        })
            except Exception as block_error:
                print(f"Error parsing block: {str(block_error)}")
                continue
                
    except Exception as e:
        print(f"Error parsing questions: {str(e)}")
    
    return questions



class MultimodalQuizView(APIView):
    def get(self, request, pk):
        try:
            # Get user by pk
            user = User.objects.get(id=pk)
            return Response({
                'message': 'Please send a POST request with content to generate quiz',
                'supported_content_types': ['text', 'image', 'audio', 'video', 'pdf', 'word', 'ppt', 'excel', 'url'],
                'difficulty_levels': ['easy', 'medium', 'hard'],
                'question_types': ['mcq', 'true_false']
            })
        except User.DoesNotExist:
            return Response({
                'error': 'User not found'
            }, status=status.HTTP_404_NOT_FOUND)

    def post(self, request, pk):
        try:
            # Get user by pk
            user = User.objects.get(id=pk)
            
            # Check for valid payment first
            # user = request.user
            # valid_payment = Payment.objects.filter(
            #     user=user,
            #     status='success',
            #     attempts_remaining__gt=0
            # ).first()
            
            # if not valid_payment:
            #     return Response({
            #         'error': 'No valid payment found. Please purchase quiz attempts.'
            #     }, status=status.HTTP_402_PAYMENT_REQUIRED)

            # # Handle multipart form data or JSON
            
            content_type = request.headers.get('Content-Type', '')
            is_multipart = 'multipart/form-data' in content_type.lower()

            if is_multipart:
                content_text = request.POST.get('content', '')
                image = request.FILES.get('image')
                audio = request.FILES.get('audio')
                video = request.FILES.get('video')
                pdf = request.FILES.get('pdf')
                word = request.FILES.get('word')
                ppt = request.FILES.get('ppt')
                excel = request.FILES.get('excel')
                url = request.POST.get('url')
                difficulty = request.POST.get('difficulty', 'medium')
                question_type = request.POST.get('question_type', 'mcq')
            else:
                content_text = request.data.get('content', '')
                image = None
                audio = None
                video = None
                pdf = None
                word = None
                ppt = None
                excel = None
                url = request.data.get('url')
                difficulty = request.data.get('difficulty', 'medium')
                question_type = request.data.get('question_type', 'mcq')

            if not any([content_text.strip(), image, audio, video, pdf, word, ppt, excel, url]):
                return Response({
                    'error': 'Please provide at least one type of content (text/image/audio/video/pdf/word/ppt/excel/url)'
                }, status=status.HTTP_400_BAD_REQUEST)

            if difficulty not in ['easy', 'medium', 'hard']:
                return Response({
                    'error': 'Invalid difficulty level. Choose from: easy, medium, hard'
                }, status=status.HTTP_400_BAD_REQUEST)

            if question_type not in ['mcq', 'true_false']:
                return Response({
                    'error': 'Invalid question type. Choose from: mcq, true_false'
                }, status=status.HTTP_400_BAD_REQUEST)

            # Generate quiz prompt based on difficulty and question type
            difficulty_instructions = {
                'easy': 'Generate basic, straightforward questions suitable for beginners.',
                'medium': 'Generate moderately challenging questions that require good understanding.',
                'hard': 'Generate complex questions that require deep understanding and critical thinking.'
            }

            if question_type == 'mcq':
                prompt = f"""
                You're an AI quiz generator.
                {difficulty_instructions[difficulty]}
                Based on the following content, generate 10 multiple choice questions with 4 options.
                Format strictly like:
                Q: <question>
                A. <option>
                B. <option>
                C. <option>
                D. <option>
                Answer: <correct_option_letter>

                Text: {content_text}
                """
            else:  # true_false
                prompt = f"""
                You're an AI quiz generator.
                {difficulty_instructions[difficulty]}
                Based on the following content, generate 10 true/false questions.
                Format strictly like:
                Q: <question>
                A. True
                B. False
                Answer: <correct_option_letter>

                Text: {content_text}
                """

            parts = [{"text": prompt}]

            # Process multimedia content
            if image:
                img_data = image.read()
                parts.append({
                    "inline_data": {
                        "mime_type": image.content_type,
                        "data": base64.b64encode(img_data).decode()
                    }
                })

            if audio:
                try:
                    transcribed = transcribe_audio(audio)
                    if not transcribed:
                        return Response({
                            'error': 'Could not transcribe audio file'
                        }, status=status.HTTP_400_BAD_REQUEST)
                    
                    content_text = transcribed
                    parts[0]["text"] = parts[0]["text"].replace("Text: ", f"Text: {content_text}")
                    parts.append({"text": f"Additional context from audio: {transcribed}"})
                    
                except Exception as e:
                    return Response({
                        'error': f'Error processing audio: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)

            if video:
                frame_data, mime = extract_frame(video)
                parts.append({
                    "inline_data": {
                        "mime_type": mime,
                        "data": base64.b64encode(frame_data).decode()
                    }
                })
                
            if pdf:
                pdf_text = extract_pdf_text(pdf)
                if pdf_text:
                    content_text = pdf_text
                    parts[0]["text"] = parts[0]["text"].replace("Text: ", f"Text: {content_text}")

            if word:
                try:
                    word_text = extract_word_text(word)
                    if word_text:
                        content_text = word_text
                        parts[0]["text"] = parts[0]["text"].replace("Text: ", f"Text: {content_text}")
                    else:
                        return Response({
                            'error': 'Could not extract text from Word document'
                        }, status=status.HTTP_400_BAD_REQUEST)
                except Exception as e:
                    return Response({
                        'error': f'Error processing Word document: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)

            if ppt:
                try:
                    ppt_text = extract_ppt_text(ppt)
                    if ppt_text:
                        content_text = ppt_text
                        parts[0]["text"] = parts[0]["text"].replace("Text: ", f"Text: {content_text}")
                    else:
                        return Response({
                            'error': 'Could not extract text from PowerPoint document'
                        }, status=status.HTTP_400_BAD_REQUEST)
                except Exception as e:
                    return Response({
                        'error': f'Error processing PowerPoint document: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)

            if excel:
                try:
                    excel_text = extract_excel_text(excel)
                    if excel_text:
                        content_text = excel_text
                        parts[0]["text"] = parts[0]["text"].replace("Text: ", f"Text: {content_text}")
                    else:
                        return Response({
                            'error': 'Could not extract text from Excel document'
                        }, status=status.HTTP_400_BAD_REQUEST)
                except Exception as e:
                    return Response({
                        'error': f'Error processing Excel document: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)
            
            if url:
                try:
                    url_text = extract_url_text(url)
                    if url_text:
                        content_text = url_text
                        parts[0]["text"] = parts[0]["text"].replace("Text:", f"Text:{content_text}")
                    else:
                        return Response({
                            'error': 'Could not extract text from URL'
                        })
                except Exception as e:
                    return Response({
                        'error': f'Error processing URL: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)

            # Generate questions
            model = genai.GenerativeModel("models/gemini-1.5-flash")
            response = model.generate_content(parts)
            output = response.text
            questions = parse_questions(output, question_type)

            if not questions:
                return Response({
                    'error': 'Failed to generate questions'
                }, status=status.HTTP_400_BAD_REQUEST)

            # Only save if answers are submitted
            if request.data.get('submitted') or (is_multipart and request.POST.get('submitted')):
                user_answers = []
                if is_multipart:
                    user_answers = [request.POST.get(f'question_{i}') for i in range(len(questions))]
                else:
                    user_answers = request.data.get('user_answers', [])
                
                score = sum(1 for i, q in enumerate(questions) if user_answers[i] == q['answer'])
                
                # Create quiz data and validate with serializer
                quiz_data = {
                    'user': user,
                    'title': f"{difficulty.capitalize()} {question_type.upper()} Quiz",
                    'questions': questions,
                    'difficulty': difficulty,
                    'question_type': question_type,
                    'content_type': 'text'
                }
                quiz_serializer = QuizSerializer(data=quiz_data)
                if not quiz_serializer.is_valid():
                    return Response(quiz_serializer.errors, status=status.HTTP_400_BAD_REQUEST)
                quiz = quiz_serializer.save()

                # Create quiz attempt data and validate with serializer
                quiz_attempt_data = {
                    'user': user,
                    'questions': questions,
                    'user_answers': user_answers,
                    'score': score,
                    'total': len(questions),
                    'difficulty': difficulty,
                    'question_type': question_type
                }
                quiz_attempt_serializer = QuizAttemptSerializer(data=quiz_attempt_data)
                if not quiz_attempt_serializer.is_valid():
                    return Response(quiz_attempt_serializer.errors, status=status.HTTP_400_BAD_REQUEST)
                quiz_attempt = quiz_attempt_serializer.save()

                return Response({
                    'quiz_id': str(quiz.id),
                    'attempt_id': str(quiz_attempt.id),
                    'score': score,
                    'total': len(questions),
                    'accuracy': (score / len(questions)) * 100 if len(questions) > 0 else 0
                })
            else:
                # If not submitted, just return the questions
                return Response({
                    'questions': [{
                        'question': q['question'],
                        'options': q['options'],
                        'answer':q['answer']
                    } for q in questions]
                })

        except User.DoesNotExist:
            return Response({
                'error': 'User not found'
            }, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            return Response({
                'error': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class UserSignupView(APIView):
    def post(self, request):
        try:
            data = request.data
            serializer = UserSerializer(data=data)
            
            if serializer.is_valid():
                # Create user instance but don't save yet
                user = User(
                    profile=data['profile'],
                    full_name=data['full_name'],
                    phone_number=data['phone_number'],
                    country_code=data['country_code'],
                    email=data['email']
                )
                
                # Set password
                user.set_password(data['password'], data['confirm_password'])
                user.save()
                
                return Response({
                    'message': 'User created successfully',
                    'user_id': str(user.id)
                }, status=status.HTTP_201_CREATED)
            
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
            
        except ValidationError as e:
            return Response({'error': str(e)}, status=status.HTTP_400_BAD_REQUEST)
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class UserLoginView(APIView):
    def post(self, request):
        try:
            email = request.data.get('email')
            password = request.data.get('password')
            
            try:
                user = User.objects.get(email=email)
            except DoesNotExist:
                return Response({'error': 'Invalid credentials'}, status=status.HTTP_401_UNAUTHORIZED)
            
            if user.check_password(password):
                user.last_login = datetime.utcnow()
                user.save()
                
                return Response({
                    'message': 'Login successful',
                    'user_id': str(user.id)
                })
            
            return Response({'error': 'Invalid credentials'}, status=status.HTTP_401_UNAUTHORIZED)
            
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class UserDetailView(APIView):
    def get(self, request, pk):
        try:
            user = User.objects.get(id=pk)
            user_data = {
                "full_name": user.full_name,
                "phone_number": user.phone_number,
                "country_code": user.country_code,
                "email": user.email,
                "is_active": user.is_active
            }
            return Response({"user": user_data})
        
        except (DoesNotExist, ValidationError):
            return Response({"error": "User not found or invalid ID"}, status=status.HTTP_404_NOT_FOUND)
    

        
    